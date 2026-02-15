# ==========================================
# Video-to-PPTX Extraction App (Queue + Logs + UI)
# ==========================================
import streamlit as st
import cv2
import tempfile
import os
import io
import uuid
from datetime import datetime
from pptx import Presentation

try:
    from scenedetect import SceneManager, open_video
    from scenedetect.detectors import ContentDetector
except ImportError:
    st.error("PySceneDetect is not installed. Please run: pip install scenedetect[opencv]")
    st.stop()

# ==========================================
# 1. PAGE CONFIG & SESSION STATE
# ==========================================
st.set_page_config(
    page_title="Video to Slides Converter",
    page_icon="🎞️",
    layout="wide",
    initial_sidebar_state="expanded",
)

def init_session():
    if "queue" not in st.session_state:
        st.session_state.queue = []
    if "logs" not in st.session_state:
        st.session_state.logs = []
    if "processing" not in st.session_state:
        st.session_state.processing = False
    if "pending_uploads" not in st.session_state:
        st.session_state.pending_uploads = []

init_session()

def log(msg: str):
    ts = datetime.now().strftime("%H:%M:%S")
    st.session_state.logs.append(f"[{ts}] {msg}")

# ==========================================
# 2. SIDEBAR: SETTINGS + UPLOAD LIMIT
# ==========================================
st.sidebar.header("⚙️ Settings")

max_upload_mb = st.sidebar.number_input(
    "Max upload size (MB)",
    min_value=10,
    max_value=2048,
    value=1000,
    step=100,
    help="Files larger than this will be rejected when added to the queue.",
)
MAX_BYTES = int(max_upload_mb * 1024 * 1024)

st.sidebar.markdown("---")
st.sidebar.subheader("Detection")

sensitivity = st.sidebar.slider(
    "Detection Sensitivity",
    min_value=0,
    max_value=100,
    value=100,
    help="Higher = detect more slide changes.",
)
threshold_val = max(5.0, 80 - (sensitivity * 0.75))

min_scene_len_sec = st.sidebar.number_input(
    "Minimum Slide Duration (sec)",
    min_value=0.2,
    max_value=10.0,
    value=0.2,
    step=0.1,
)

st.sidebar.markdown("---")
st.sidebar.subheader("Deduplication")
dedup_slides = st.sidebar.checkbox(
    "Merge slides that differ only by person",
    value=True,
    help="Mask faces and keep one slide when content is the same (reduces duplicates from professor movement).",
)
similarity_threshold = st.sidebar.slider(
    "Similarity threshold (merge if above)",
    min_value=0.80,
    max_value=0.99,
    value=0.92,
    step=0.01,
    help="Higher = only merge very similar slides. Lower = merge more aggressively.",
    disabled=not dedup_slides,
)

# ==========================================
# 3. PROCESSING ENGINE
# ==========================================
_FACE_CASCADE = None

def _get_face_cascade():
    global _FACE_CASCADE
    if _FACE_CASCADE is None:
        path = os.path.join(cv2.data.haarcascades, "haarcascade_frontalface_default.xml")
        _FACE_CASCADE = cv2.CascadeClassifier(path)
    return _FACE_CASCADE


def _mask_faces(frame_bgr, padding_ratio=0.2):
    """Mask face regions with blurred background so slide content can be compared."""
    gray = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2GRAY)
    cascade = _get_face_cascade()
    faces = cascade.detectMultiScale(gray, scaleFactor=1.1, minNeighbors=5, minSize=(30, 30))
    out = frame_bgr.copy()
    for (x, y, w, h) in faces:
        pad_w = int(w * padding_ratio)
        pad_h = int(h * padding_ratio)
        x1 = max(0, x - pad_w)
        y1 = max(0, y - pad_h)
        x2 = min(frame_bgr.shape[1], x + w + pad_w)
        y2 = min(frame_bgr.shape[0], y + h + pad_h)
        roi = out[y1:y2, x1:x2]
        if roi.size > 0:
            blurred = cv2.GaussianBlur(roi, (0, 0), 30)
            out[y1:y2, x1:x2] = blurred
    return out


def _slide_similarity(img1_bgr, img2_bgr):
    """Compare two images (after face masking) using histogram correlation. Returns 0-1, higher = more similar."""
    if img1_bgr is None or img2_bgr is None:
        return 0.0
    # Resize to same size for fair comparison
    h, w = 200, 200
    a = cv2.resize(img1_bgr, (w, h))
    b = cv2.resize(img2_bgr, (w, h))
    # HSV histograms (2D H-S) are good for content comparison
    hsv_a = cv2.cvtColor(a, cv2.COLOR_BGR2HSV)
    hsv_b = cv2.cvtColor(b, cv2.COLOR_BGR2HSV)
    hist_a = cv2.calcHist([hsv_a], [0, 1], None, [50, 60], [0, 180, 0, 256])
    hist_b = cv2.calcHist([hsv_b], [0, 1], None, [50, 60], [0, 180, 0, 256])
    cv2.normalize(hist_a, hist_a, 0, 1, cv2.NORM_MINMAX)
    cv2.normalize(hist_b, hist_b, 0, 1, cv2.NORM_MINMAX)
    corr = cv2.compareHist(hist_a, hist_b, cv2.HISTCMP_CORREL)
    # correlation is in [-1, 1]; map to [0, 1]
    return (corr + 1) / 2.0


def detect_slides(video_path, threshold, min_duration):
    video = open_video(video_path)
    min_scene_len_frames = int(min_duration * video.frame_rate)
    scene_manager = SceneManager()
    detector = ContentDetector(threshold=threshold, min_scene_len=min_scene_len_frames)
    scene_manager.add_detector(detector)
    scene_manager.detect_scenes(video, show_progress=False)
    return scene_manager.get_scene_list(), video


def generate_pptx(video_path, scene_list, progress_callback=None, dedup_slides=False, similarity_threshold=0.92):
    if not scene_list:
        return None, 0
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        return None, 0
    pptx_buffer = io.BytesIO()
    total_scenes = len(scene_list)
    prev_masked = None
    slides_added = 0

    for i, scene in enumerate(scene_list):
        if progress_callback:
            progress_callback(i + 1, total_scenes, f"Extracting frame {i+1}/{total_scenes}")
        # scene is (start_timecode, end_timecode) tuple; start_timecode has get_frames()
        start_timecode = scene[0] if isinstance(scene, tuple) else scene
        start_frame = start_timecode.get_frames() + 10
        cap.set(cv2.CAP_PROP_POS_FRAMES, start_frame)
        success, frame = cap.read()
        if not success:
            continue
        # Face-aware dedup: if only person differs, skip this slide
        if dedup_slides:
            masked = _mask_faces(frame)
            if prev_masked is not None:
                sim = _slide_similarity(masked, prev_masked)
                if sim >= similarity_threshold:
                    continue  # same slide content, skip duplicate
            prev_masked = masked
        frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        ok, buffer = cv2.imencode(".jpg", frame_rgb)
        if ok:
            image_stream = io.BytesIO(buffer)
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                image_stream,
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            slides_added += 1
    cap.release()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer, slides_added


# ==========================================
# 4. MAIN LAYOUT: HEADER + UPLOAD + QUEUE
# ==========================================
st.title("🎞️ Video to PowerPoint Converter")
st.markdown("Extract slides from video using **scene detection**. Add multiple files to the queue and process them one by one.")
st.markdown("---")

col_upload, col_queue = st.columns([1, 1])

with col_upload:
    st.subheader("📤 Upload videos")
    st.caption("Choose files → click **Open**. After upload finishes, the list below will appear; then click **Add to queue**.")
    uploaded = st.file_uploader(
        "Select one or more videos (MP4, MOV, AVI, MKV)",
        type=["mp4", "mov", "avi", "mkv"],
        accept_multiple_files=True,
        key="video_upload",
    )

    # Persist uploaded files into session state so they survive reruns (e.g. after "Add to queue" click)
    if uploaded:
        st.session_state.pending_uploads = []
        for u in uploaded:
            try:
                data = u.read()
                if len(data) > MAX_BYTES:
                    log(f"Rejected {u.name}: over {max_upload_mb} MB limit")
                    st.session_state.pending_uploads.append({
                        "name": u.name, "size": len(data), "data": None, "rejected": True,
                    })
                else:
                    st.session_state.pending_uploads.append({
                        "name": u.name, "size": len(data), "data": data, "rejected": False,
                    })
            except Exception as e:
                st.session_state.pending_uploads.append({
                    "name": u.name, "size": 0, "data": None, "rejected": True, "error": str(e),
                })

    # Show selected files so user sees immediate feedback after "Open"
    pending = st.session_state.pending_uploads
    if pending:
        st.success(f"**{len([p for p in pending if not p.get('rejected')])} file(s) selected** — click **Add to queue** below.")
        for p in pending:
            if p.get("rejected"):
                st.warning(f"⚠️ {p['name']} — over size limit or error (max {max_upload_mb} MB)")
            else:
                size_mb = p["size"] / (1024 * 1024)
                st.caption(f"📄 {p['name']} ({size_mb:.2f} MB)")
        st.markdown("")

    def add_to_queue():
        for p in st.session_state.pending_uploads:
            if p.get("rejected") or p.get("data") is None:
                continue
            st.session_state.queue.append({
                "id": str(uuid.uuid4()),
                "name": p["name"],
                "size": p["size"],
                "data": p["data"],
                "status": "pending",
                "pptx_data": None,
                "error": None,
                "slide_count": None,
            })
            log(f"Added to queue: {p['name']} ({p['size'] / (1024*1024):.1f} MB)")
        added = len([p for p in st.session_state.pending_uploads if not p.get("rejected") and p.get("data")])
        st.session_state.pending_uploads = []
        if added:
            st.toast(f"Added {added} file(s) to queue", icon="✅")

    if pending and any(not p.get("rejected") and p.get("data") for p in pending):
        st.button("Add to queue", type="primary", on_click=add_to_queue, key="add_to_queue_btn")

with col_queue:
    st.subheader("📋 Queue")
    if not st.session_state.queue:
        st.info("Queue is empty. Upload files and click **Add to queue**.")
    else:
        for i, item in enumerate(st.session_state.queue):
            status_icon = {"pending": "⏳", "processing": "🔄", "done": "✅", "error": "❌"}.get(item["status"], "⏳")
            with st.container():
                st.markdown(f"{status_icon} **{item['name']}** — {item['status']}")
                if item["slide_count"] is not None:
                    st.caption(f"Slides: {item['slide_count']}")
                if item["error"]:
                    st.caption(f"Error: {item['error']}")
            st.markdown("")

# ==========================================
# 5. PROCESS QUEUE BUTTON + STATUS BOX + LOGS
# ==========================================
st.markdown("---")
run_col, _ = st.columns([1, 3])
with run_col:
    process_btn = st.button("🚀 Process queue", type="primary", use_container_width=True)

pending = [q for q in st.session_state.queue if q["status"] == "pending"]
if process_btn and pending and not st.session_state.processing:
    st.session_state.processing = True

if st.session_state.processing and pending:
    status_placeholder = st.empty()
    progress_placeholder = st.progress(0)
    total_to_process = len(pending)
    current_file_index = 0

    for item in st.session_state.queue:
        if item["status"] != "pending":
            continue

        item["status"] = "processing"
        current_file_index += 1
        total_in_queue = total_to_process
        status_placeholder.markdown(f"""
        <div style="background: var(--secondary-background-color); padding: 1rem; border-radius: 8px; border-left: 4px solid #6366f1;">
        <strong>🔄 Current</strong>: File {current_file_index} of {total_in_queue}<br>
        <strong>File</strong>: {item['name']}<br>
        <strong>Step</strong>: <span id="step-text">Detecting scenes...</span>
        </div>
        """, unsafe_allow_html=True)
        log(f"Processing: {item['name']}")

        tfile = tempfile.NamedTemporaryFile(delete=False, suffix=".mp4")
        tfile.write(item["data"])
        temp_path = tfile.name
        tfile.close()

        try:
            scenes, _ = detect_slides(temp_path, threshold_val, min_scene_len_sec)
            status_placeholder.markdown(f"""
            <div style="background: var(--secondary-background-color); padding: 1rem; border-radius: 8px; border-left: 4px solid #6366f1;">
            <strong>🔄 Current</strong>: File {current_file_index} of {total_in_queue}<br>
            <strong>File</strong>: {item['name']}<br>
            <strong>Step</strong>: Building PPTX ({len(scenes)} slides)...
            </div>
            """, unsafe_allow_html=True)

            def update_step(current, total, msg):
                progress_placeholder.progress(current / total)
                status_placeholder.markdown(f"""
                <div style="background: var(--secondary-background-color); padding: 1rem; border-radius: 8px; border-left: 4px solid #6366f1;">
                <strong>🔄 Current</strong>: File {current_file_index} of {total_in_queue}<br>
                <strong>File</strong>: {item['name']}<br>
                <strong>Step</strong>: {msg}
                </div>
                """, unsafe_allow_html=True)

            pptx_data, num_slides = generate_pptx(
                temp_path, scenes,
                progress_callback=update_step,
                dedup_slides=dedup_slides,
                similarity_threshold=similarity_threshold,
            )
            if pptx_data:
                item["pptx_data"] = pptx_data.getvalue()
                item["slide_count"] = num_slides
                item["status"] = "done"
                log(f"Done: {item['name']} — {num_slides} slides")
            else:
                item["status"] = "error"
                item["error"] = "No slides detected or PPTX build failed"
                log(f"Error: {item['name']} — {item['error']}")
        except Exception as e:
            item["status"] = "error"
            item["error"] = str(e)
            log(f"Error: {item['name']} — {e}")
        finally:
            try:
                os.unlink(temp_path)
            except Exception:
                pass

    st.session_state.processing = False
    progress_placeholder.progress(1.0)
    st.rerun()

# ==========================================
# 6. LIVE STATUS BOX (when not processing)
# ==========================================
if not st.session_state.processing and st.session_state.queue:
    processing_items = [q for q in st.session_state.queue if q["status"] == "processing"]
    if processing_items:
        st.info("Processing in progress...")
    else:
        done = [q for q in st.session_state.queue if q["status"] == "done"]
        if done:
            st.success(f"✅ {len(done)} file(s) ready to download below.")

# ==========================================
# 7. LOGS PANEL
# ==========================================
st.markdown("---")
with st.expander("📜 Logs", expanded=False):
    if st.session_state.logs:
        log_text = "\n".join(reversed(st.session_state.logs[-100:]))
        st.text_area("", value=log_text, height=180, disabled=True, label_visibility="collapsed")
        if st.button("Clear logs"):
            st.session_state.logs = []
            st.rerun()
    else:
        st.caption("No logs yet.")

# ==========================================
# 8. DOWNLOADS: COMPLETED FILES
# ==========================================
st.markdown("---")
st.subheader("📥 Download results")

done_items = [q for q in st.session_state.queue if q["status"] == "done" and q.get("pptx_data")]
if not done_items:
    st.caption("Completed files will appear here with a download button.")
else:
    for item in done_items:
        base = os.path.splitext(item["name"])[0]
        fname = f"Slides_{base}_{datetime.now().strftime('%Y-%m-%d_%H%M')}.pptx"
        st.download_button(
            label=f"📥 Download {item['name']} ({item['slide_count']} slides)",
            data=bytes(item["pptx_data"]),
            file_name=fname,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=item["id"],
        )

# ==========================================
# 9. CLEAR QUEUE
# ==========================================
st.sidebar.markdown("---")
if st.sidebar.button("Clear queue & logs"):
    st.session_state.queue = []
    st.session_state.logs = []
    st.rerun()

st.sidebar.caption("Supported: MP4, MOV, AVI, MKV")
